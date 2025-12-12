import os
import io
from typing import Optional, Tuple
from werkzeug.utils import secure_filename
from werkzeug.datastructures import FileStorage

try:
    from PyPDF2 import PdfReader
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from PIL import Image
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class FileService:
    """Service for file upload, storage, and processing"""

    ALLOWED_EXTENSIONS = {
        'pdf': ['pdf'],
        'image': ['png', 'jpg', 'jpeg', 'gif', 'bmp'],
        'ppt': ['ppt', 'pptx'],
        'audio': ['mp3', 'wav', 'm4a'],
        'video': ['mp4', 'mov', 'avi']
    }

    def __init__(self, upload_folder: str = None):
        self.upload_folder = upload_folder or os.getenv(
            'UPLOAD_FOLDER',
            os.path.join(os.path.dirname(os.path.dirname(__file__)), 'uploads')
        )
        # Create upload folder if it doesn't exist
        os.makedirs(self.upload_folder, exist_ok=True)
        os.makedirs(os.path.join(self.upload_folder, 'documents'), exist_ok=True)
        os.makedirs(os.path.join(self.upload_folder, 'images'), exist_ok=True)
        os.makedirs(os.path.join(self.upload_folder, 'thumbnails'), exist_ok=True)

    def allowed_file(self, filename: str, file_type: str = None) -> bool:
        """Check if file extension is allowed"""
        if '.' not in filename:
            return False

        ext = filename.rsplit('.', 1)[1].lower()

        if file_type:
            return ext in self.ALLOWED_EXTENSIONS.get(file_type, [])

        # Check if extension is in any category
        return any(ext in extensions for extensions in self.ALLOWED_EXTENSIONS.values())

    def save_file(self, file: FileStorage, subfolder: str = 'documents') -> Tuple[str, str, int]:
        """
        Save uploaded file and return (filename, file_path, file_size)
        """
        if not file or file.filename == '':
            raise ValueError("No file provided")

        if not self.allowed_file(file.filename):
            raise ValueError(f"File type not allowed: {file.filename}")

        # Secure the filename
        filename = secure_filename(file.filename)

        # Add timestamp to avoid collisions
        from datetime import datetime
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        name, ext = os.path.splitext(filename)
        filename = f"{name}_{timestamp}{ext}"

        # Save file
        folder_path = os.path.join(self.upload_folder, subfolder)
        os.makedirs(folder_path, exist_ok=True)
        file_path = os.path.join(folder_path, filename)

        file.save(file_path)

        # Get file size
        file_size = os.path.getsize(file_path)

        return filename, file_path, file_size

    def extract_text_from_pdf(self, file_path: str) -> Tuple[Optional[str], int]:
        """
        Extract text from PDF file
        Returns (extracted_text, page_count)
        """
        if not PYPDF2_AVAILABLE:
            return None, 0

        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                page_count = len(pdf_reader.pages)

                text_content = []
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)

                extracted_text = "\n\n".join(text_content)
                return extracted_text, page_count

        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {str(e)}")

    def extract_text_from_image(self, file_path: str) -> Optional[str]:
        """
        Extract text from image using OCR
        """
        if not TESSERACT_AVAILABLE:
            return None

        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            raise Exception(f"Error performing OCR: {str(e)}")

    def extract_text_from_ppt(self, file_path: str) -> Tuple[Optional[str], int]:
        """
        Extract text from PowerPoint file
        Returns (extracted_text, slide_count)
        """
        if not PPTX_AVAILABLE:
            return None, 0

        try:
            prs = Presentation(file_path)
            slide_count = len(prs.slides)

            text_content = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)

                if slide_text:
                    text_content.append("\n".join(slide_text))

            extracted_text = "\n\n".join(text_content)
            return extracted_text, slide_count

        except Exception as e:
            raise Exception(f"Error extracting text from PPT: {str(e)}")

    def generate_thumbnail(self, file_path: str, output_path: str = None,
                          size: Tuple[int, int] = (200, 200)) -> Optional[str]:
        """
        Generate thumbnail for image file
        """
        if not TESSERACT_AVAILABLE:  # PIL is imported with pytesseract
            return None

        try:
            image = Image.open(file_path)

            # Convert to RGB if necessary
            if image.mode in ('RGBA', 'LA', 'P'):
                image = image.convert('RGB')

            # Create thumbnail
            image.thumbnail(size, Image.Resampling.LANCZOS)

            # Save thumbnail
            if not output_path:
                name, _ = os.path.splitext(os.path.basename(file_path))
                output_path = os.path.join(
                    self.upload_folder, 'thumbnails', f"{name}_thumb.jpg"
                )

            image.save(output_path, 'JPEG')
            return output_path

        except Exception as e:
            raise Exception(f"Error generating thumbnail: {str(e)}")

    def process_document(self, file_path: str, file_type: str) -> Dict:
        """
        Process uploaded document (extract text, generate thumbnail, etc.)
        Returns dict with processing results
        """
        result = {
            'extracted_text': None,
            'ocr_text': None,
            'page_count': 0,
            'thumbnail_url': None
        }

        try:
            if file_type == 'PDF' and PYPDF2_AVAILABLE:
                text, pages = self.extract_text_from_pdf(file_path)
                result['extracted_text'] = text
                result['page_count'] = pages

            elif file_type == 'PPT' and PPTX_AVAILABLE:
                text, slides = self.extract_text_from_ppt(file_path)
                result['extracted_text'] = text
                result['page_count'] = slides

            elif file_type == 'IMAGE' and TESSERACT_AVAILABLE:
                ocr_text = self.extract_text_from_image(file_path)
                result['ocr_text'] = ocr_text
                # Generate thumbnail for images
                thumb_path = self.generate_thumbnail(file_path)
                if thumb_path:
                    result['thumbnail_url'] = os.path.basename(thumb_path)

        except Exception as e:
            print(f"Error processing document: {str(e)}")

        return result

    def delete_file(self, file_path: str) -> bool:
        """Delete file from filesystem"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                return True
            return False
        except Exception as e:
            raise Exception(f"Error deleting file: {str(e)}")


# Singleton instance
_file_service = None

def get_file_service() -> FileService:
    """Get or create file service instance"""
    global _file_service
    if _file_service is None:
        _file_service = FileService()
    return _file_service
